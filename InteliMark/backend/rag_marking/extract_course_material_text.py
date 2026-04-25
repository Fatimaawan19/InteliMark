import json
import os
import re
import sys
import time
import tempfile
import subprocess
from pathlib import Path


def _build_ocr_engine(enable_ocr: bool):
    """Return OCR callable and diagnostic metadata.

    OCR is intentionally optional because Tesseract binary may not be installed
    on every machine. We fail soft and keep the text pipeline usable.
    """

    if not enable_ocr:
        return None, "disabled", ""

    try:
        import io
        from PIL import Image  # type: ignore[import-not-found]
        import pytesseract  # type: ignore[import-not-found]

        tesseract_cmd = os.environ.get("TESSERACT_CMD", "").strip()
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

        # Validate Tesseract binary availability early for clear diagnostics.
        _ = pytesseract.get_tesseract_version()

        def _ocr_from_bytes(image_bytes: bytes) -> str:
            img = Image.open(io.BytesIO(image_bytes))
            return (pytesseract.image_to_string(img) or "").strip()

        return _ocr_from_bytes, "ready", ""
    except Exception as ocr_error:
        return None, "unavailable", str(ocr_error)


def _extract_with_pymupdf(file_path: str, ocr_reader=None):
    import fitz  # PyMuPDF  # type: ignore[import-not-found]

    doc = fitz.open(file_path)
    pages = []
    full_text = []
    total_images = 0
    ocr_blocks = []
    max_ocr_images = 20
    ocr_images_seen = 0

    for i, page in enumerate(doc, start=1):  # type: ignore
        text = page.get_text("text") or ""
        images = page.get_images(full=True)
        image_count = len(images)
        total_images += image_count
        pages.append({"page": i, "char_count": len(text)})
        full_text.append(text)

        if ocr_reader and ocr_images_seen < max_ocr_images:
            for img_info in images[:3]:
                if ocr_images_seen >= max_ocr_images:
                    break
                try:
                    xref = img_info[0]
                    image_bytes = doc.extract_image(xref).get("image")
                    if not image_bytes:
                        continue
                    ocr_text = ocr_reader(image_bytes)
                    if ocr_text:
                        ocr_blocks.append(
                            {
                                "page": i,
                                "text": ocr_text[:1200],
                            }
                        )
                    ocr_images_seen += 1
                except Exception:
                    continue
    combined = "\n".join(full_text).strip()
    return {
        "text": combined,
        "page_count": len(doc),
        "pages": pages,
        "extractor": "pymupdf",
        "visual_count": total_images,
        "source_format": "pdf",
        "ocr_blocks": ocr_blocks,
    }


def _extract_with_pdfplumber(file_path: str):
    import pdfplumber  # type: ignore[import-not-found]

    pages = []
    full_text = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append({"page": i, "char_count": len(text)})
            full_text.append(text)
    combined = "\n".join(full_text).strip()
    return {
        "text": combined,
        "page_count": len(pages),
        "pages": pages,
        "extractor": "pdfplumber",
        "visual_count": 0,
        "source_format": "pdf",
        "ocr_blocks": [],
    }


def _extract_from_pptx(file_path: str, ocr_reader=None):
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation(file_path)
    slides_meta = []
    all_text = []
    total_visuals = 0
    ocr_blocks = []
    max_ocr_images = 20
    ocr_images_seen = 0

    for idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        visual_count = 0

        for shape in slide.shapes:
            # text boxes / titles
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                txt = (getattr(shape, "text", "") or "").strip()
                if txt:
                    slide_text_parts.append(txt)

            # Only count actual images (shapes with .image attribute)
            shape_type = str(getattr(shape, "shape_type", ""))
            has_image = hasattr(shape, "image") and getattr(shape, "image", None) is not None

            if has_image:
                try:
                    image_obj = getattr(shape, "image", None)
                    if image_obj is None:
                        continue
                    image_bytes = image_obj.blob
                    if not image_bytes:
                        continue
                    
                    # We have actual image data - count it
                    visual_count += 1
                    
                    # Try to extract OCR if enabled and within limit
                    if ocr_reader and ocr_images_seen < max_ocr_images:
                        try:
                            ocr_text = ocr_reader(image_bytes).strip()
                            ocr_blocks.append({
                                "page": idx,
                                "text": ocr_text[:1200] if ocr_text else "[Image processed - no text detected]"
                            })
                            ocr_images_seen += 1
                        except Exception:
                            pass
                except Exception:
                    # Shape claimed to have image but failed to access - skip it
                    continue

        total_visuals += visual_count

        slide_text = "\n".join(slide_text_parts).strip()
        slides_meta.append(
            {
                "page": idx,
                "char_count": len(slide_text),
                "visual_count": visual_count,
            }
        )

        all_text.append(f"[SLIDE {idx}]\n{slide_text}" if slide_text else f"[SLIDE {idx}]")

    combined = "\n\n".join(all_text).strip()
    return {
        "text": combined,
        "page_count": len(prs.slides),
        "pages": slides_meta,
        "extractor": "python-pptx",
        "visual_count": total_visuals,
        "source_format": "pptx",
        "ocr_blocks": ocr_blocks,
    }


def _convert_ppt_to_pptx(file_path: str) -> str:
    """Convert legacy .ppt to .pptx using LibreOffice (soffice).

    Requires LibreOffice installed and `soffice` available in PATH.
    Returns path to converted pptx in a temp folder.
    """
    input_path = Path(file_path)
    if input_path.suffix.lower() != ".ppt":
        return str(input_path)

    soffice = os.environ.get("SOFFICE_BIN", "").strip() or "soffice"
    out_dir = Path(tempfile.mkdtemp(prefix="ppt-convert-"))

    # LibreOffice CLI: soffice --headless --convert-to pptx --outdir <dir> <file>
    cmd = [
        soffice,
        "--headless",
        "--nologo",
        "--nolockcheck",
        "--norestore",
        "--convert-to",
        "pptx",
        "--outdir",
        str(out_dir),
        str(input_path),
    ]

    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    except FileNotFoundError as e:
        raise RuntimeError(
            "LibreOffice is required to extract .ppt files. Install LibreOffice and ensure `soffice` is on PATH "
            "(or set SOFFICE_BIN to full path)."
        ) from e
    except subprocess.TimeoutExpired as e:
        raise RuntimeError("Timed out converting .ppt to .pptx using LibreOffice.") from e

    if proc.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed (code {proc.returncode}). stderr: {proc.stderr.strip()[:400]}"
        )

    converted = out_dir / (input_path.stem + ".pptx")
    if not converted.exists():
        # LibreOffice sometimes preserves original extension casing; fallback to search.
        matches = list(out_dir.glob("*.pptx"))
        if matches:
            converted = matches[0]
        else:
            raise RuntimeError("LibreOffice conversion did not produce a .pptx output.")

    return str(converted)


def _extract_from_docx(file_path: str):
    from docx import Document  # type: ignore[import-not-found]

    doc = Document(file_path)
    parts = []
    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        if txt:
            parts.append(txt)

    # Extract basic table text too (helps many lecture notes).
    for table in doc.tables:
        for row in table.rows:
            cells = []
            for cell in row.cells:
                t = (cell.text or "").strip()
                if t:
                    cells.append(t)
            if cells:
                parts.append(" | ".join(cells))

    combined = "\n".join(parts).strip()
    return {
        "text": combined,
        "page_count": 1,
        "pages": [{"page": 1, "char_count": len(combined), "visual_count": 0}],
        "extractor": "python-docx",
        "visual_count": 0,
        "source_format": "docx",
        "ocr_blocks": [],
    }


_EQUATION_PATTERNS = [
    r"\b[a-zA-Z]\s*=\s*[^\n]+",  # x = ...
    r"[0-9a-zA-Z]+\s*[+\-*/^]\s*[0-9a-zA-Z]+",  # algebraic expressions
    r"\b(integral|derivative|sigma|theta|lambda|alpha|beta)\b",
]


def _extract_equation_candidates(text: str):
    if not text:
        return []

    candidates = []
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    for ln in lines:
        for pattern in _EQUATION_PATTERNS:
            if re.search(pattern, ln, flags=re.IGNORECASE):
                candidates.append(ln)
                break

    # keep unique and short list
    unique = []
    seen = set()
    for item in candidates:
        key = item.lower()
        if key not in seen:
            seen.add(key)
            unique.append(item)
    return unique[:50]


def _build_augmented_text(base_text: str, equation_candidates, visual_count: int, ocr_blocks):
    parts = [base_text.strip()] if base_text.strip() else []

    if equation_candidates:
        parts.append("[EQUATION_CANDIDATES]\n" + "\n".join(f"- {e}" for e in equation_candidates))

    if visual_count > 0:
        parts.append(
            "[VISUAL_HINTS]\n"
            f"- This material contains approximately {visual_count} visual elements (images/charts/diagrams).\n"
            "- If user asks about a diagram/flowchart, mention that visual interpretation may need OCR/vision parsing for full fidelity."
        )

    if ocr_blocks:
        image_lines = []
        for block in ocr_blocks[:25]:
            txt = block.get("text", "").strip()
            if not txt:
                continue
            image_lines.append(f"- page {block.get('page', '?')}: {txt}")
        if image_lines:
            parts.append("[IMAGE_TEXT]\n" + "\n".join(image_lines))

    return "\n\n".join(parts).strip()


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "Missing file path argument", "reason": "No file path was provided to the extractor."}))
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(json.dumps({"success": False, "error": f"File not found: {file_path}", "reason": "The uploaded file path does not exist on disk."}))
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()
    enable_ocr = str(os.environ.get("RAG_ENABLE_OCR", "false")).lower() in {"1", "true", "yes", "on"}
    ocr_reader, ocr_status, ocr_error = _build_ocr_engine(enable_ocr)

    started = time.time()
    try:
        if ext == ".pdf":
            try:
                result = _extract_with_pymupdf(file_path, ocr_reader=ocr_reader)
            except Exception:
                result = _extract_with_pdfplumber(file_path)
        elif ext in {".pptx", ".ppt"}:
            pptx_path = _convert_ppt_to_pptx(file_path) if ext == ".ppt" else file_path
            result = _extract_from_pptx(pptx_path, ocr_reader=ocr_reader)
            if ext == ".ppt":
                result["source_format"] = "ppt"
                result["extractor"] = "libreoffice+python-pptx"
        elif ext == ".docx":
            result = _extract_from_docx(file_path)
        else:
            print(
                json.dumps(
                    {
                        "success": False,
                        "error": f"Unsupported file type: {ext}",
                        "reason": "Only PDF, PPT/PPTX, and DOCX files are supported by the extractor.",
                    }
                )
            )
            sys.exit(1)
    except Exception as extraction_error:
        reason_text = str(extraction_error).strip()
        if not reason_text:
            reason_text = repr(extraction_error)
        print(
            json.dumps(
                {
                    "success": False,
                    "error": "Extraction failed",
                    "reason": reason_text,
                    "error_type": type(extraction_error).__name__,
                }
            )
        )
        sys.exit(1)

    duration_ms = int((time.time() - started) * 1000)
    base_text = result["text"]
    equation_candidates = _extract_equation_candidates(base_text)
    text = _build_augmented_text(
        base_text,
        equation_candidates=equation_candidates,
        visual_count=int(result.get("visual_count", 0) or 0),
        ocr_blocks=result.get("ocr_blocks", []) or [],
    )

    ocr_blocks = result.get("ocr_blocks", []) or []
    ocr_char_count = sum(len((b.get("text") or "")) for b in ocr_blocks)

    print(
        json.dumps(
            {
                "success": True,
                "text": text,
                "char_count": len(text),
                "page_count": result["page_count"],
                "pages": result["pages"],
                "extractor": result["extractor"],
                "equation_count": len(equation_candidates),
                "visual_count": int(result.get("visual_count", 0) or 0),
                "ocr_count": len(ocr_blocks),
                "ocr_char_count": ocr_char_count,
                "ocr_blocks": ocr_blocks,
                "ocr_enabled": bool(ocr_reader),
                "ocr_status": ocr_status,
                "ocr_error": ocr_error,
                "source_format": result.get("source_format", ext.lstrip(".")),
                "duration_ms": duration_ms,
            }
        )
    )


if __name__ == "__main__":
    main()
